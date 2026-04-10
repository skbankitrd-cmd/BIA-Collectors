from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import uuid
from datetime import datetime

class PPTExporter:
    """原子化 PPT 導出器：全系統統一投影片規範"""
    
    def __init__(self):
        self.output_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../../BIA-Web-Solution/server/public/generated_reports"))
        os.makedirs(self.output_dir, exist_ok=True)
        
        # 金控配色規範
        self.primary_color = RGBColor(230, 0, 18) # 台新紅

    def export(self, slides_data: list, report_title: str) -> str:
        """
        將資料轉為 PPT
        slides_data 格式: [{"title": "標題", "content": "內容"}]
        """
        prs = Presentation()
        
        # 建立首頁
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = report_title
        subtitle.text = f"產出日期: {datetime.now().strftime('%Y-%m-%d %H:%M')}\nBIA 戰略情報系統"
        
        # 設定標題顏色
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color

        # 建立內容頁
        bullet_slide_layout = prs.slide_layouts[1]
        for data in slides_data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = data['title']
            slide.placeholders[1].text = data['content']

        # 儲存
        file_name = f"PPT_{report_title}_{uuid.uuid4().hex[:4]}.pptx"
        full_path = os.path.join(self.output_dir, file_name)
        prs.save(full_path)
        
        return f"/public/generated_reports/{file_name}"
