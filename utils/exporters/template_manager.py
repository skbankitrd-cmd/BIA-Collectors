import os
import openpyxl
from pptx import Presentation
import logging
from datetime import datetime
import uuid

logger = logging.getLogger("TemplateManager")

class TemplateManager:
    """智慧模板管理器：支援 Excel 與 PPT 範本填充"""
    
    def __init__(self):
        self.base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../../"))
        self.template_path = os.path.join(self.base_dir, "templates")
        self.output_dir = os.path.join(self.base_dir, "BIA-Web-Solution/server/public/generated_reports")
        os.makedirs(self.template_path, exist_ok=True)
        os.makedirs(self.output_dir, exist_ok=True)

    def fill_excel(self, template_name, data_map):
        # [Harness Security Fix: Path Traversal Prevention]
        safe_name = os.path.basename(template_name)
        wb = openpyxl.load_workbook(os.path.join(self.template_path, safe_name))
        ws = wb.active
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and "{{" in cell.value:
                    key = cell.value.replace("{{", "").replace("}}", "").strip()
                    if key in data_map: cell.value = data_map[key]
        
        filename = f"Filled_{uuid.uuid4().hex[:4]}.xlsx"
        wb.save(os.path.join(self.output_dir, filename))
        return f"/public/generated_reports/{filename}"

    def fill_ppt(self, template_name, data_map):
        """讀取 PPT 範本並替換標籤"""
        # [Harness Security Fix: Path Traversal Prevention]
        safe_name = os.path.basename(template_name)
        prs = Presentation(os.path.join(self.template_path, safe_name))
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{" in run.text:
                            for key, val in data_map.items():
                                run.text = run.text.replace(f"{{{{{key}}}}}", str(val))
        
        filename = f"Filled_{uuid.uuid4().hex[:4]}.pptx"
        prs.save(os.path.join(self.output_dir, filename))
        return f"/public/generated_reports/{filename}"

    def fill_any(self, template_name: str, data_map: dict) -> str:
        if template_name.endswith('.xlsx'):
            return self.fill_excel(template_name, data_map)
        elif template_name.endswith('.pptx'):
            return self.fill_ppt(template_name, data_map)
        return ""
